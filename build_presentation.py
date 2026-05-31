"""Generate QuizGen.pptx — pitch deck for the QuizGen project.

Структура строго следует «Требованиям к материалам для хакатона»
(Сбер Образование / Школа 21):
  1. Название прототипа
  2. Команда (название, ФИО, роли и зоны ответственности)
  3. Проблема и пользователи (формулировка, целевая группа, сценарий)
  4. Решение (идея, способ достижения цели, ценность)
  5. Основные пользовательские сценарии (пользовательский путь)
  6. Сравнение с существующими решениями
  7. Техническая реализация — архитектура (клиент ↔ сервер ↔ LLM ↔ БД)
  8. Техническая реализация — поток данных и принципы работы
  9. Стек технологий
 10. MVP (минимальная версия, ресурсы, возможности развития)
 11. Дорожная карта (таблица)
 12. Риски и меры нейтрализации (таблица)
 13. Блок дополнительной информации + ссылки
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─── Palette ─────────────────────────────────────────────────────────────────
BLUE_DARK  = RGBColor(0x1E, 0x3A, 0x8A)
BLUE       = RGBColor(0x25, 0x63, 0xEB)
BLUE_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)
BLUE_BG    = RGBColor(0xEF, 0xF6, 0xFF)
SLATE_900  = RGBColor(0x0F, 0x17, 0x2A)
SLATE_700  = RGBColor(0x33, 0x41, 0x55)
SLATE_500  = RGBColor(0x64, 0x74, 0x8B)
SLATE_300  = RGBColor(0xCB, 0xD5, 0xE1)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN      = RGBColor(0x16, 0xA3, 0x4A)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)
RED        = RGBColor(0xDC, 0x26, 0x26)

FONT = "Segoe UI"
MONO = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.shadow.inherit = False
    return s


def add_rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.0)
    shp.shadow.inherit = False
    return shp


def _split_bold(s):
    """'foo **bar** baz' → [('foo ',False),('bar',True),(' baz',False)]"""
    parts, buf, bold = [], "", False
    i = 0
    while i < len(s):
        if s[i:i + 2] == "**":
            if buf:
                parts.append((buf, bold))
                buf = ""
            bold = not bold
            i += 2
        else:
            buf += s[i]
            i += 1
    if buf:
        parts.append((buf, bold))
    return parts or [("", False)]


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=SLATE_700,
             font=FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for seg, seg_bold in _split_bold(line):
            r = p.add_run()
            r.text = seg
            r.font.name = font
            r.font.size = Pt(size)
            r.font.bold = bold or seg_bold
            r.font.color.rgb = SLATE_900 if (seg_bold and not bold) else color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=18, color=SLATE_700,
                bullet_color=BLUE, line_spacing=1.3, marker="▍"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        br = p.add_run()
        br.text = f"{marker}  "
        br.font.name = FONT
        br.font.size = Pt(size)
        br.font.bold = True
        br.font.color.rgb = bullet_color
        for seg, is_bold in _split_bold(item):
            r = p.add_run()
            r.text = seg
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = is_bold
            r.font.color.rgb = SLATE_900 if is_bold else color
    return tb


def header(slide, kicker, title, subtitle=None):
    """kicker = номер/секция шаблона, title = заголовок слайда."""
    add_rect(slide, Inches(0.6), Inches(0.5), Inches(0.18), Inches(0.62), BLUE)
    if kicker:
        add_text(slide, Inches(0.95), Inches(0.42), Inches(11.5), Inches(0.32),
                 kicker.upper(), size=12, bold=True, color=BLUE)
    add_text(slide, Inches(0.95), Inches(0.72), Inches(11.5), Inches(0.6),
             title, size=29, bold=True, color=BLUE_DARK)
    if subtitle:
        add_text(slide, Inches(0.95), Inches(1.32), Inches(11.5), Inches(0.4),
                 subtitle, size=15, color=SLATE_500)
    add_rect(slide, Inches(0.6), Inches(1.78), Inches(12.1), Emu(12700), SLATE_300)


def footer(slide, page=None, total=None):
    add_text(slide, Inches(0.6), Inches(7.12), Inches(8), Inches(0.3),
             "QuizGen · СберХакатон 2026", size=10, color=SLATE_500)
    if page is not None:
        add_text(slide, Inches(11.5), Inches(7.12), Inches(1.5), Inches(0.3),
                 f"{page} / {total}", size=10, color=SLATE_500,
                 align=PP_ALIGN.RIGHT)


def tag(slide, x, y, text, *, fill=BLUE_LIGHT, color=BLUE_DARK, size=12):
    w = Inches(0.105 * len(text) + 0.4)
    h = Inches(0.35)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    return x + w + Inches(0.1)


def box(slide, x, y, w, h, title, body=None, *, fill=BLUE_BG, title_color=BLUE_DARK,
        body_color=SLATE_700, title_size=15, body_size=12, line=None, accent=None):
    add_rect(slide, x, y, w, h, fill, line=line)
    if accent:
        add_rect(slide, x, y, w, Inches(0.08), accent)
    ty = y + Inches(0.18)
    add_text(slide, x + Inches(0.2), ty, w - Inches(0.4), Inches(0.5),
             title, size=title_size, bold=True, color=title_color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if body:
        add_text(slide, x + Inches(0.2), ty + Inches(0.5), w - Inches(0.4),
                 h - Inches(0.75),
                 body, size=body_size, color=body_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.25)


def arrow(slide, x, y, w, h, *, shape=MSO_SHAPE.RIGHT_ARROW, fill=BLUE):
    a = add_rect(slide, x, y, w, h, fill, shape=shape)
    return a


def styled_table(slide, x, y, w, h, headers, rows, col_widths=None, *,
                 header_fill=BLUE, fsize=13, hsize=14):
    tbl_shape = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h)
    tbl = tbl_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    for i, htxt in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text_frame.clear()
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = htxt
        r.font.name = FONT
        r.font.size = Pt(hsize)
        r.font.bold = True
        r.font.color.rgb = WHITE
        cell.margin_left = Inches(0.15)
        cell.margin_right = Inches(0.1)
        cell.margin_top = Inches(0.06)
        cell.margin_bottom = Inches(0.06)
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 else BLUE_BG
            cell.text_frame.clear()
            cell.text_frame.word_wrap = True
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            for seg, is_bold in _split_bold(str(val)):
                r = p.add_run()
                r.text = seg
                r.font.name = FONT
                r.font.size = Pt(fsize)
                r.font.bold = is_bold
                r.font.color.rgb = SLATE_900 if is_bold else SLATE_700
            cell.margin_left = Inches(0.15)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
    return tbl


def screen_mock(slide, x, y, w, h, name, lines, *, accent=BLUE):
    """Лёгкий wireframe «окна интерфейса» с заголовком и условным содержимым."""
    add_rect(slide, x, y, w, h, WHITE, line=SLATE_300)
    bar_h = Inches(0.34)
    add_rect(slide, x, y, w, bar_h, accent)
    add_text(slide, x + Inches(0.18), y, w - Inches(0.3), bar_h,
             name, size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(0.18), y + bar_h + Inches(0.1),
             w - Inches(0.34), h - bar_h - Inches(0.2),
             "\n".join(lines), size=10.5, color=SLATE_700, line_spacing=1.25)


# ═══════════════════════════════════════════════════════════════════════════════
# 1. TITLE — Название прототипа
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_rect(s, 0, 0, Inches(4.4), SH, BLUE_DARK)
add_rect(s, Inches(4.4), 0, Inches(0.12), SH, BLUE)

add_rect(s, Inches(0.8), Inches(0.9), Inches(2.8), Inches(2.8), BLUE)
add_text(s, Inches(0.8), Inches(0.9), Inches(2.8), Inches(2.8),
         "QG", size=120, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.8), Inches(4.05), Inches(3.2), Inches(0.4),
         "СБЕРХАКАТОН 2026", size=13, bold=True, color=BLUE_LIGHT)
add_text(s, Inches(0.8), Inches(4.45), Inches(3.2), Inches(0.4),
         "ИИ для образования: автоматизация рутины", size=12, color=BLUE_LIGHT)

add_text(s, Inches(5.0), Inches(1.7), Inches(8), Inches(1.2),
         "QuizGen", size=72, bold=True, color=BLUE_DARK)
add_text(s, Inches(5.0), Inches(2.8), Inches(8), Inches(0.6),
         "AI-генератор викторин для учителей", size=24, color=SLATE_700)

add_rect(s, Inches(5.0), Inches(3.7), Inches(0.7), Inches(0.06), BLUE)

add_text(s, Inches(5.0), Inches(4.0), Inches(8), Inches(1.6),
         "Загружаешь учебный материал —\nGigaChat генерирует квиз.\n"
         "Ученики проходят его по ссылке.\nУчитель получает аналитику.",
         size=18, color=SLATE_500, line_spacing=1.3)

x = Inches(5.0)
for t in ["Go 1.22", "React 18", "PostgreSQL", "GigaChat", "Docker"]:
    x = tag(s, x, Inches(6.0), t, size=13)

footer(s)

# ═══════════════════════════════════════════════════════════════════════════════
# 2. КОМАНДА
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header(s, "Шаг 2 · Команда", "Команда «QuizGen»",
       "ФИО участников, распределение ролей и зон ответственности")

styled_table(s, Inches(0.95), Inches(2.05), Inches(11.4), Inches(3.9),
             ["Участник (ФИО)", "Роль", "Зона ответственности"],
             [
                 ["**Фамилия Имя**",   "Backend / Tech Lead",
                  "Go-сервис, REST API, интеграция с GigaChat, БД, Docker"],
                 ["**Фамилия Имя**",   "Frontend",
                  "React SPA, UI/UX, редактор квиза, плеер ученика, графики"],
                 ["**Фамилия Имя**",   "ML / Промпт-инжиниринг",
                  "Промпты, парсинг PDF/DOCX/PPTX, качество генерации"],
                 ["**Фамилия Имя**",   "Product / Дизайн",
                  "Продуктовое исследование, карточка прототипа, презентация"],
             ],
             col_widths=[Inches(3.0), Inches(3.2), Inches(5.2)])

cx, cy, cw, ch = Inches(0.95), Inches(6.15), Inches(11.4), Inches(0.75)
add_rect(s, cx, cy, cw, ch, BLUE_BG)
add_rect(s, cx, cy, Inches(0.12), ch, AMBER)
add_text(s, cx + Inches(0.35), cy, cw - Inches(0.5), ch,
         "⚠ Замените «Фамилия Имя» на реальные ФИО участников команды.",
         size=13, bold=True, color=BLUE_DARK, anchor=MSO_ANCHOR.MIDDLE)

footer(s, 2, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 3. ПРОБЛЕМА И ПОЛЬЗОВАТЕЛИ
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header(s, "Шаг 3 · Проблема и пользователи", "Учитель тонет в рутине проверки знаний",
       "Формулировка проблемы · целевая группа · сценарий столкновения")

# Проблема (левый блок)
add_text(s, Inches(0.95), Inches(2.0), Inches(6.0), Inches(0.4),
         "Проблема", size=16, bold=True, color=BLUE_DARK)
add_bullets(s, Inches(0.95), Inches(2.5), Inches(6.0), Inches(3),
            [
                "Составление проверочной по новому материалу — **1–2 часа** ручной работы.",
                "Готовые сервисы требуют **ручного ввода** вопросов и ответов.",
                "Нет связки «материал → квиз → аналитика класса» в один клик.",
                "Раздать индивидуально и защитить от списывания — **отдельная боль**.",
            ], size=14, line_spacing=1.35)

# Целевая группа (правый блок)
add_text(s, Inches(7.25), Inches(2.0), Inches(5.1), Inches(0.4),
         "Целевая группа", size=16, bold=True, color=BLUE_DARK)
for i, (t, d) in enumerate([
        ("Учитель школы", "основной пользователь: создаёт и раздаёт квизы"),
        ("Преподаватель СПО/вуза", "быстрые срезы знаний по лекции"),
        ("Репетитор / тьютор", "проверка после занятия, домашние тесты"),
        ("Ученик", "проходит квиз по ссылке/QR с телефона")]):
    yy = Inches(2.5 + i * 0.72)
    add_rect(s, Inches(7.25), yy, Inches(0.1), Inches(0.6), BLUE)
    add_text(s, Inches(7.5), yy, Inches(4.9), Inches(0.3),
             t, size=13.5, bold=True, color=SLATE_900)
    add_text(s, Inches(7.5), yy + Inches(0.3), Inches(4.9), Inches(0.3),
             d, size=12, color=SLATE_500)

# Сценарий столкновения (нижняя плашка)
sy = Inches(5.7)
add_rect(s, Inches(0.95), sy, Inches(11.4), Inches(1.1), BLUE_DARK)
add_text(s, Inches(1.25), sy + Inches(0.12), Inches(11), Inches(0.35),
         "Сценарий столкновения с проблемой", size=13, bold=True, color=BLUE_LIGHT)
add_text(s, Inches(1.25), sy + Inches(0.45), Inches(11), Inches(0.55),
         "Вечер. Учителю нужен тест по новой главе на завтра. Он вручную выдумывает "
         "вопросы, верстает их, печатает варианты, потом вручную проверяет 30 работ — "
         "и всё равно не видит, какую тему провалил класс.",
         size=13.5, color=WHITE, line_spacing=1.25)

footer(s, 3, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 4. РЕШЕНИЕ
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header(s, "Шаг 4 · Решение", "QuizGen закрывает весь цикл",
       "Идея сервиса · решение как способ достижения цели")

# Идея — одна фраза
iy = Inches(1.95)
add_rect(s, Inches(0.95), iy, Inches(11.4), Inches(0.95), BLUE)
add_text(s, Inches(1.25), iy, Inches(10.9), Inches(0.95),
         "Загрузил материал или задал тему → GigaChat сгенерировал квиз → "
         "раздал ученикам по ссылке/QR → получил аналитику с разбором.",
         size=16, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)

# Две колонки ценности
cy1, cw1, ch1 = Inches(3.15), Inches(5.7), Inches(3.0)
box(s, Inches(0.95), cy1, cw1, ch1, "Для учителя", accent=BLUE, fill=BLUE_BG)
add_bullets(s, Inches(1.25), cy1 + Inches(0.7), cw1 - Inches(0.6), ch1 - Inches(0.9),
            [
                "Генерация из **PDF / DOCX / PPTX / TXT / MD**",
                "Настройка: сложность, тон, уровень Блума, типы",
                "Редактор + **перегенерация одного вопроса**",
                "Персональные ссылки и **QR**, печать с ключом",
                "Аналитика по классу и **CSV-экспорт**",
            ], size=14, line_spacing=1.3)

cx2 = Inches(7.05)
box(s, cx2, cy1, cw1, ch1, "Для ученика", accent=BLUE_DARK, fill=BLUE_BG)
add_bullets(s, cx2 + Inches(0.3), cy1 + Inches(0.7), cw1 - Inches(0.6), ch1 - Inches(0.9),
            [
                "Открыть ссылку с телефона или **QR**",
                "Представиться по имени",
                "Пройти квиз с **прогресс-баром**",
                "Сразу увидеть процент правильных",
                "Честно: **антишпаргалка** работает в фоне",
            ], size=14, line_spacing=1.3, bullet_color=BLUE_DARK)

x = Inches(0.95)
for t in ["~30 сек на квиз", "5 форматов", "QR-ссылки", "Антишпаргалка", "CSV"]:
    x = tag(s, x, Inches(6.45), t, size=12)

footer(s, 4, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 5. ОСНОВНЫЕ ПОЛЬЗОВАТЕЛЬСКИЕ СЦЕНАРИИ (пользовательский путь)
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header(s, "Шаг 4 · Пользовательский путь", "Полный сценарий: от материала до статистики",
       "Пошаговый путь по экранам интерфейса")

screens = [
    ("1 · Генерация", ["Тема / класс / сложность",
                       "Прикрепить файл", "[ Сгенерировать ]"]),
    ("2 · Редактор",  ["Список вопросов",
                       "Правка текста и ответов", "↻ Перегенерировать"]),
    ("3 · Раздача",   ["Publish → ссылки",
                       "QR на каждого", "Печать с ключом"]),
    ("4 · Плеер",     ["Вход по имени",
                       "Вопросы + прогресс", "Антишпаргалка"]),
    ("5 · Аналитика", ["Гистограмма баллов",
                       "Таблица учеников", "Экспорт CSV"]),
]
n = len(screens)
gap = Inches(0.2)
total_w = Inches(12.2)
card_w = Emu(int((total_w - gap * (n - 1)) / n))
card_h = Inches(2.6)
y0 = Inches(2.3)
x0 = Inches(0.57)
for i, (name, lines) in enumerate(screens):
    x = x0 + Emu(int(card_w + gap) * i)
    accent = BLUE_DARK if i in (3,) else BLUE
    screen_mock(s, x, y0, card_w, card_h, name, lines, accent=accent)
    if i < n - 1:
        ax = x + card_w + Emu(int(gap)) // 4
        arrow(s, ax, y0 + Inches(1.0), Inches(0.16), Inches(0.5),
              shape=MSO_SHAPE.RIGHT_ARROW, fill=SLATE_300)

# подпись пути
py = Inches(5.4)
add_rect(s, Inches(0.95), py, Inches(11.4), Inches(1.05), BLUE_BG)
add_rect(s, Inches(0.95), py, Inches(0.12), Inches(1.05), BLUE)
add_text(s, Inches(1.3), py, Inches(10.9), Inches(1.05),
         "Воспроизводимость полного сценария: ввод контекста и загрузка → генерация → "
         "редактирование → воспроизведение по ссылке → экспорт и статистика — "
         "без выхода из приложения.",
         size=14, color=SLATE_700, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)

footer(s, 5, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 6. СРАВНЕНИЕ С СУЩЕСТВУЮЩИМИ РЕШЕНИЯМИ
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header(s, "Шаг 4 · Сравнение", "Чем QuizGen отличается от аналогов",
       "Сравнение существующих решений и предлагаемого сервиса")

styled_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(4.6),
             ["Возможность", "Kahoot /\nQuizizz", "Quizlet", "Google\nForms", "**QuizGen**"],
             [
                 ["Авто-генерация из материала (PDF/DOCX/PPTX)", "—", "—", "—", "**✓ GigaChat**"],
                 ["Ввод вопросов вручную не требуется",          "—", "—", "—", "**✓**"],
                 ["Персональные ссылки + QR на ученика",         "частично", "—", "ручные", "**✓ авто**"],
                 ["Антишпаргалка (вкладки, копипаст)",           "—", "—", "—", "**✓**"],
                 ["Аналитика с разбором попытки + CSV",          "базовая", "базовая", "базовая", "**детальная**"],
                 ["Настройка по таксономии Блума и сложности",   "—", "—", "—", "**✓**"],
                 ["Self-hosted, данные у школы (Docker)",        "—", "—", "—", "**✓**"],
             ],
             col_widths=[Inches(4.6), Inches(1.85), Inches(1.6), Inches(1.65), Inches(2.4)],
             fsize=12.5, hsize=12.5)

footer(s, 6, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 7. ТЕХНИЧЕСКАЯ РЕАЛИЗАЦИЯ — АРХИТЕКТУРА
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header(s, "Шаг 5 · Техническая реализация", "Архитектура: клиент ↔ сервер ↔ LLM ↔ БД",
       "Слоистый монолит на Go, фронт — отдельный SPA")

# Browser
box(s, Inches(0.7), Inches(3.0), Inches(2.5), Inches(1.3),
    "Браузер", "React SPA\n(учитель + ученик)", fill=BLUE_BG, accent=BLUE,
    title_size=15, body_size=12)
# HTTPS arrow
arrow(s, Inches(3.3), Inches(3.45), Inches(0.85), Inches(0.45),
      shape=MSO_SHAPE.LEFT_RIGHT_ARROW, fill=BLUE)
add_text(s, Inches(3.25), Inches(3.12), Inches(1.0), Inches(0.3),
         "HTTPS", size=10, bold=True, color=SLATE_500, align=PP_ALIGN.CENTER)

# Server box with 3 layers
srv_x, srv_y, srv_w, srv_h = Inches(4.35), Inches(2.35), Inches(4.0), Inches(2.95)
add_rect(s, srv_x, srv_y, srv_w, srv_h, WHITE, line=BLUE)
add_text(s, srv_x, srv_y + Inches(0.1), srv_w, Inches(0.35),
         "Go (Gin) :8080", size=14, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
for i, (lab) in enumerate(["Handlers — HTTP-роутинг",
                           "Services — бизнес-логика + LLM",
                           "Repositories — доступ к БД"]):
    ly = srv_y + Inches(0.6 + i * 0.72)
    add_rect(s, srv_x + Inches(0.3), ly, srv_w - Inches(0.6), Inches(0.58), BLUE_BG)
    add_text(s, srv_x + Inches(0.3), ly, srv_w - Inches(0.6), Inches(0.58),
             lab, size=12, bold=True, color=SLATE_700,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 2:
        arrow(s, srv_x + srv_w / 2 - Inches(0.1), ly + Inches(0.58), Inches(0.2), Inches(0.14),
              shape=MSO_SHAPE.DOWN_ARROW, fill=SLATE_300)

# GigaChat
arrow(s, Inches(8.5), Inches(2.95), Inches(0.8), Inches(0.4),
      shape=MSO_SHAPE.RIGHT_ARROW, fill=BLUE)
add_text(s, Inches(8.45), Inches(2.62), Inches(0.95), Inches(0.3),
         "OAuth/REST", size=9, bold=True, color=SLATE_500, align=PP_ALIGN.CENTER)
box(s, Inches(9.45), Inches(2.65), Inches(3.0), Inches(1.05),
    "GigaChat", "LLM API (Sber)\nгенерация квиза", fill=BLUE_DARK,
    title_color=WHITE, body_color=BLUE_LIGHT, title_size=15, body_size=11)

# PostgreSQL
arrow(s, Inches(8.5), Inches(4.45), Inches(0.8), Inches(0.4),
      shape=MSO_SHAPE.RIGHT_ARROW, fill=BLUE)
add_text(s, Inches(8.55), Inches(4.12), Inches(0.9), Inches(0.3),
         "SQL", size=9, bold=True, color=SLATE_500, align=PP_ALIGN.CENTER)
box(s, Inches(9.45), Inches(4.15), Inches(3.0), Inches(1.05),
    "PostgreSQL :5432", "квизы · сессии\nответы · статистика", fill=BLUE_BG,
    accent=BLUE, title_size=14, body_size=11)

add_text(s, Inches(0.95), Inches(5.75), Inches(11.4), Inches(1.0),
         "Принцип: чистое разделение слоёв — HTTP-роутинг, бизнес-логика и доступ к данным "
         "развязаны, меняются и тестируются независимо. Провайдер LLM абстрагирован: "
         "GigaChat по умолчанию, fallback на Anthropic/OpenAI одной переменной окружения.",
         size=13, color=SLATE_500, line_spacing=1.3)

footer(s, 7, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 8. ТЕХ. РЕАЛИЗАЦИЯ — ПОТОК ДАННЫХ И ПРИНЦИПЫ
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header(s, "Шаг 5 · Поток данных", "Как материал превращается в квиз и статистику",
       "Диаграмма потока данных · интеграция с LLM · хранение")

flow = ["Файл /\nтема", "Парсер\nPDF·DOCX·PPTX", "Сборка\nпромпта", "GigaChat",
        "JSON-парсинг\n+ sanitizeUTF8", "PostgreSQL", "Ссылки\n+ QR", "Статистика"]
n = len(flow)
gap = Inches(0.12)
fw = Emu(int((Inches(12.2) - gap * (n - 1)) / n))
fy = Inches(2.05)
fx0 = Inches(0.57)
for i, lab in enumerate(flow):
    x = fx0 + Emu(int(fw + gap) * i)
    fill = BLUE_DARK if lab == "GigaChat" else BLUE_BG
    tc = WHITE if lab == "GigaChat" else SLATE_900
    add_rect(s, x, fy, fw, Inches(1.0), fill, line=(None if lab == "GigaChat" else SLATE_300))
    add_text(s, x, fy, fw, Inches(1.0), lab, size=10.5, bold=True, color=tc,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    if i < n - 1:
        add_text(s, x + fw - Emu(int(gap)) // 2, fy, Emu(int(gap)) + Emu(20000), Inches(1.0),
                 "›", size=16, bold=True, color=SLATE_500,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_bullets(s, Inches(0.95), Inches(3.55), Inches(11.4), Inches(3.2),
            [
                "**Интеграция с LLM (GigaChat):** OAuth-токен по Authorization-Key кешируется до expires_at; вызов через REST.",
                "**Промпты:** жёсткая JSON-схема, правила для каждого типа вопроса (single / multiple / true-false), уровни сложности и таксономия Блума, учёт типичных ошибок; отдельный промпт для перегенерации одного вопроса.",
                "**Парсинг форматов:** PDF (pdfcpu), DOCX/PPTX (zip + XML), TXT/MD напрямую; лимит 10 МБ, обрезка до 6000 символов, санитайзинг UTF-8/NUL.",
                "**Хранение данных:** PostgreSQL — квизы, вопросы, сессии (session ID), ответы, метаданные, статистика и история версий.",
                "**Лимиты и экспорт:** rate-limit генераций к LLM на пользователя; вывод — персональные ссылки с session ID, полноэкранный плеер, печать, экспорт CSV/PDF.",
            ], size=13, line_spacing=1.3)

footer(s, 8, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 9. СТЕК ТЕХНОЛОГИЙ
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header(s, "Шаг 5 · Стек", "Технологический стек",
       "Современный, проверенный, ничего лишнего")

rows = [
    ["**Backend**",  "Go 1.22 · Gin · golang-migrate · lib/pq"],
    ["**Database**", "PostgreSQL 16"],
    ["**LLM**",      "GigaChat (Sber) — OAuth + REST · fallback: Anthropic / OpenAI"],
    ["**Frontend**", "React 18 · TypeScript · Vite · TailwindCSS · React Router"],
    ["**Charts**",   "Recharts (гистограмма баллов)"],
    ["**QR**",       "qrcode.react"],
    ["**Парсинг**",  "pdfcpu · archive/zip + XML (DOCX / PPTX)"],
    ["**Авторизация**", "JWT (HMAC-SHA256) · middleware Auth + RateLimit"],
    ["**DevOps**",   "Docker · docker-compose (multi-stage: node → go → alpine)"],
]
styled_table(s, Inches(0.95), Inches(2.0), Inches(11.4), Inches(4.7),
             ["Слой", "Технологии"], rows,
             col_widths=[Inches(2.6), Inches(8.8)])

footer(s, 9, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 10. MVP
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header(s, "Шаг 6 · MVP", "Минимальная рабочая версия — уже работает",
       "Что сделано · ресурсы · возможности развития")

done = [
    "Регистрация и логин учителя (JWT)",
    "Генерация квиза из текста или файла",
    "Настройка: сложность, тон, Блум, типы",
    "Редактор + перегенерация вопроса",
    "Персональные ссылки + QR-коды",
    "Печать — с ключом ответов и без",
    "Прохождение учеником + антишпаргалка",
    "Статистика: гистограмма, таблица, разбор",
    "CSV-экспорт (Excel-friendly)",
    "Docker-развёртка одной командой",
]
add_text(s, Inches(0.95), Inches(1.95), Inches(7.0), Inches(0.35),
         "Минимальная версия (готово)", size=15, bold=True, color=BLUE_DARK)
half = (len(done) + 1) // 2
for col, items in enumerate([done[:half], done[half:]]):
    cx = Inches(0.95 + col * 3.65)
    for i, item in enumerate(items):
        y = Inches(2.45 + i * 0.52)
        add_rect(s, cx, y + Inches(0.03), Inches(0.28), Inches(0.28), GREEN,
                 shape=MSO_SHAPE.OVAL)
        add_text(s, cx, y + Inches(0.03), Inches(0.28), Inches(0.28),
                 "✓", size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx + Inches(0.4), y, Inches(3.1), Inches(0.45),
                 item, size=12, color=SLATE_700, anchor=MSO_ANCHOR.MIDDLE)

# Ресурсы
rx, ry, rw = Inches(8.55), Inches(1.95), Inches(3.85)
add_text(s, rx, ry, rw, Inches(0.35), "Ресурсы для MVP", size=15, bold=True, color=BLUE_DARK)
add_bullets(s, rx, ry + Inches(0.5), rw, Inches(2.0),
            [
                "1 backend-инстанс Go + PostgreSQL",
                "Ключ API **GigaChat** (scope PERS)",
                "Любой Docker-хост (≈512 МБ RAM)",
                "Команда из 4 человек",
            ], size=12.5, line_spacing=1.3)

# Развитие
add_text(s, rx, Inches(4.5), rw, Inches(0.35),
         "Возможности развития", size=15, bold=True, color=BLUE_DARK)
add_bullets(s, rx, Inches(5.0), rw, Inches(1.8),
            [
                "Адаптивная сложность вопросов",
                "Банк вопросов и классы/группы",
                "Динамика по ученику, SSO, PWA",
            ], size=12.5, line_spacing=1.3, bullet_color=AMBER)

footer(s, 10, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 11. ДОРОЖНАЯ КАРТА
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header(s, "Шаг 7 · Дорожная карта", "От MVP до продукта",
       "Процесс реализации от начала до результата")

styled_table(s, Inches(0.95), Inches(2.05), Inches(11.4), Inches(4.4),
             ["Этап", "Срок", "Результат"],
             [
                 ["**MVP (хакатон)**", "Готово",
                  "Генерация, раздача по ссылкам/QR, аналитика, антишпаргалка, Docker"],
                 ["**Пилот**", "+1 месяц",
                  "Обкатка в 1 школе, банк вопросов, классы и группы, исправление багов"],
                 ["**Бета**", "+3 месяца",
                  "Адаптивная сложность, динамика по ученику, SSO школьных ID"],
                 ["**Релиз 1.0**", "+6 месяцев",
                  "PWA / мобильное приложение ученика, маркетплейс шаблонов квизов"],
             ],
             col_widths=[Inches(2.4), Inches(1.8), Inches(7.2)])

# мини-таймлайн
ty = Inches(6.55)
add_rect(s, Inches(1.2), ty, Inches(10.9), Inches(0.05), BLUE_LIGHT)
for i, lab in enumerate(["MVP", "Пилот", "Бета", "Релиз"]):
    px = Inches(1.5 + i * 3.5)
    add_rect(s, px, ty - Inches(0.07), Inches(0.18), Inches(0.18), BLUE, shape=MSO_SHAPE.OVAL)
    add_text(s, px - Inches(0.4), ty + Inches(0.15), Inches(1.0), Inches(0.3),
             lab, size=11, bold=True, color=SLATE_500, align=PP_ALIGN.CENTER)

footer(s, 11, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 12. РИСКИ
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header(s, "Шаг 8 · Риски", "Риски реализации и меры нейтрализации",
       "Что может пойти не так и как мы это закрываем")

styled_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(4.7),
             ["Риск", "Влияние", "Мера нейтрализации"],
             [
                 ["Невалидный JSON от LLM", "Высокое",
                  "Жёсткая схема промпта, strip ```json``` и висячих запятых, валидация ответа"],
                 ["Лимиты и стоимость API GigaChat", "Среднее",
                  "Rate-limit на пользователя, кеш OAuth-токена, абстракция провайдера (fallback)"],
                 ["Смешение данных между учителями", "Высокое",
                  "JWT + проверка пользователя в БД, изоляция сессий, приватные ссылки"],
                 ["Списывание учениками", "Среднее",
                  "Антишпаргалка: детект вкладок, блок копипаста, счётчик нарушений на бэке"],
                 ["Плохой парсинг файла", "Среднее",
                  "5 форматов + санитайзинг; всегда можно поправить вручную в редакторе"],
                 ["Слабое методическое качество", "Среднее",
                  "Таксономия Блума, уровни сложности, перегенерация отдельного вопроса"],
             ],
             col_widths=[Inches(3.4), Inches(1.5), Inches(7.2)], fsize=12.5)

footer(s, 12, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 13. ДОП. ИНФОРМАЦИЯ + ССЫЛКИ + СПАСИБО
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
header(s, "Доп. информация", "Артефакты и ссылки",
       "Чем хочется поделиться")

# Ссылки (кликабельные — подставьте реальные URL)
links = [
    ("🔗 Демо MVP",        "http://localhost:8080  (или ссылка на развёрнутый сервис)"),
    ("💻 Репозиторий",     "github.com/<team>/quizgen  — README с инструкцией по запуску"),
    ("🗂 Карточка прототипа", "ссылка из кейса — заполнена полностью"),
    ("🎬 Скринкаст",       "запись полного сценария (5–7 минут)"),
]
add_text(s, Inches(0.95), Inches(1.95), Inches(7.0), Inches(0.35),
         "Ссылки на артефакты", size=15, bold=True, color=BLUE_DARK)
for i, (t, d) in enumerate(links):
    y = Inches(2.5 + i * 0.82)
    add_rect(s, Inches(0.95), y, Inches(0.12), Inches(0.6), BLUE)
    add_text(s, Inches(1.25), y, Inches(6.5), Inches(0.32),
             t, size=14, bold=True, color=SLATE_900)
    add_text(s, Inches(1.25), y + Inches(0.32), Inches(6.5), Inches(0.32),
             d, size=11.5, color=SLATE_500)

# Доп. факты (правая колонка)
add_text(s, Inches(8.0), Inches(1.95), Inches(4.4), Inches(0.35),
         "Что отметить", size=15, bold=True, color=BLUE_DARK)
add_bullets(s, Inches(8.0), Inches(2.5), Inches(4.4), Inches(3.0),
            [
                "Контейнер сам прогоняет миграции при старте",
                "Данные изолированы, не доступны публично",
                "Лицензия — **MIT**",
                "Готов к пилоту в школе",
            ], size=13, line_spacing=1.35)

# Плашка «Спасибо»
ty = Inches(5.95)
add_rect(s, Inches(0.95), ty, Inches(11.4), Inches(0.95), BLUE_DARK)
add_text(s, Inches(0.95), ty + Inches(0.12), Inches(11.4), Inches(0.4),
         "Спасибо!  QuizGen — викторины за минуту, аналитика за секунду.",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.95), ty + Inches(0.52), Inches(11.4), Inches(0.35),
         "docker compose up -d   →   http://localhost:8080",
         size=12, font=MONO, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

footer(s, 13, 13)

# ─── Save ────────────────────────────────────────────────────────────────────
out = "QuizGen.pptx"
prs.save(out)
print(f"Saved: {out} ({len(prs.slides)} slides)")
